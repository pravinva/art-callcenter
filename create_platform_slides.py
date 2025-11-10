#!/usr/bin/env python3
"""
Script to create PowerPoint slides for Enterprise Agentic Knowledge Platform
Adds 4 new slides (B, C, D, and E) after the existing "Enterprise Agentic Knowledge Platform" slide
Version 2.0 - Handles different PowerPoint templates and layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def find_best_layout(prs, preferred_type="title_and_content"):
    """
    Find the best available slide layout from the presentation
    """
    print(f"\nAvailable slide layouts ({len(prs.slide_layouts)} total):")

    # Map layout types to typical indices
    layout_preferences = {
        "title": [0, 1],
        "title_and_content": [1, 2, 3, 5],
        "blank": [6, len(prs.slide_layouts) - 1, 5, 4]
    }

    for i, layout in enumerate(prs.slide_layouts):
        try:
            name = layout.name
            print(f"  [{i}] {name}")
        except:
            print(f"  [{i}] <unnamed layout>")

    # Try to find appropriate layout
    if preferred_type == "title_and_content":
        # Try common indices for title and content layouts
        for idx in layout_preferences["title_and_content"]:
            if idx < len(prs.slide_layouts):
                return prs.slide_layouts[idx]
    elif preferred_type == "blank":
        # Try common indices for blank layouts
        for idx in layout_preferences["blank"]:
            if idx < len(prs.slide_layouts):
                return prs.slide_layouts[idx]

    # Fallback to first available layout
    print(f"  → Using layout [0] as fallback")
    return prs.slide_layouts[0]

def create_text_box_with_title(slide, title_text, top=0.5):
    """
    Create a title text box manually
    """
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.LEFT
    return title_box

def create_slides(input_file="ART_Enterprise_AI_Strategy.pptx", output_file="ART_Enterprise_AI_Strategy_Enhanced.pptx"):
    """
    Creates enhanced presentation with new Enterprise Agentic Knowledge Platform slides
    """

    # Load existing presentation
    print(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)

    # Find the index where to insert new slides
    insert_index = None
    print(f"\nSearching through {len(prs.slides)} existing slides...")

    for i, slide in enumerate(prs.slides):
        # Check if slide has a title
        title_text = ""
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Check first text box or title
                    if "Enterprise Agentic Knowledge Platform" in shape.text or "Agentic Knowledge Platform" in shape.text:
                        title_text = shape.text
                        insert_index = i + 1
                        print(f"  ✓ Found target slide at position {i}: '{title_text[:50]}...'")
                        break
            if insert_index:
                break

    if insert_index is None:
        print("  ! Could not find 'Enterprise Agentic Knowledge Platform' slide.")
        print("  → Adding new slides at end of presentation")
        insert_index = len(prs.slides)

    # Find best layouts
    content_layout = find_best_layout(prs, "title_and_content")
    blank_layout = find_best_layout(prs, "blank")

    print(f"\n{'='*70}")
    print(f"Creating 4 new slides at position {insert_index}")
    print(f"{'='*70}\n")

    # ========================================
    # SLIDE B: Technical Architecture
    # ========================================
    print("Creating Slide B: Federated Multi-Agent Architecture...")
    slide_b = prs.slides.add_slide(blank_layout)

    # Add title
    create_text_box_with_title(slide_b, "Federated Multi-Agent Architecture", top=0.4)

    # Add subtitle
    subtitle_box = slide_b.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(9), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "How the Platform Works - Technical Architecture"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.italic = True
    subtitle_para.font.color.rgb = RGBColor(68, 68, 68)

    # Add main content
    content_box = slide_b.shapes.add_textbox(
        Inches(0.5), Inches(1.7), Inches(9), Inches(5.1)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    # Shared Foundation
    p = tf.paragraphs[0]
    p.text = "Shared Foundation (The Hub):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_after = Pt(6)

    bullets = [
        "Unity Catalog Functions Library: 50-100+ reusable tools",
        "Core tools: get_member_360, search_knowledge_base, check_compliance, trigger_workflow",
        "Single Vector Search index spanning all organizational knowledge",
        "Unified Model Serving infrastructure (Claude, Llama)",
        "Central governance and audit framework"
    ]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(3)

    # BU-Specific Agents
    p = tf.add_paragraph()
    p.text = "BU-Specific Agents (The Spokes):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_after = Pt(6)
    p.space_before = Pt(10)

    bullets = [
        "Each business unit deploys specialized agents",
        "Agents inherit from shared base configuration",
        "Access authorized subset of tools via Unity Catalog permissions",
        "All leverage same foundational infrastructure"
    ]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(3)

    # The Power
    p = tf.add_paragraph()
    p.text = "The Power of This Approach:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_before = Pt(10)
    p.space_after = Pt(6)

    p = tf.add_paragraph()
    p.text = "Member Services builds 18 UC Functions → Investment reuses 12, adds 6 → Compliance reuses 14, adds 4"
    p.level = 1
    p.font.size = Pt(11)
    p.space_after = Pt(3)

    p = tf.add_paragraph()
    p.text = "Each new agent costs 70% less because infrastructure exists"
    p.level = 1
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)

    # Key Technical Point
    p = tf.add_paragraph()
    p.text = "Key: UC Functions = SQL-based tools → Auditable, Governed, Testable, Versionable"
    p.font.size = Pt(11)
    p.font.italic = True
    p.space_before = Pt(10)
    p.font.color.rgb = RGBColor(68, 68, 68)

    print("  ✓ Slide B created")

    # ========================================
    # SLIDE C: Real-World Agent Examples
    # ========================================
    print("Creating Slide C: Multi-BU Agent Deployment...")
    slide_c = prs.slides.add_slide(blank_layout)

    # Add title
    create_text_box_with_title(slide_c, "Multi-BU Agent Deployment - Concrete Use Cases", top=0.4)

    # Add content in two columns
    left_width = 4.5
    right_width = 4.5
    column_height = 5.3

    # Left column
    left_box = slide_c.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(left_width), Inches(column_height)
    )
    tf_left = left_box.text_frame
    tf_left.word_wrap = True

    # Member Services Agent
    p = tf_left.paragraphs[0]
    p.text = "Member Services Agent"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)

    p = tf_left.add_paragraph()
    p.text = "Use cases: Live call assist, email automation, member portal chatbot"
    p.font.size = Pt(10)
    p.level = 1
    p.space_after = Pt(4)

    p = tf_left.add_paragraph()
    p.text = "Tools: get_member_360, search_knowledge_base, check_compliance"
    p.font.size = Pt(9)
    p.level = 1
    p.font.italic = True
    p.space_after = Pt(4)

    p = tf_left.add_paragraph()
    p.text = "Outcome: 25% faster resolution, 24/7 availability"
    p.font.size = Pt(10)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_after = Pt(12)

    # Investment Research Agent
    p = tf_left.add_paragraph()
    p.text = "Investment Research Agent"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_before = Pt(8)

    p = tf_left.add_paragraph()
    p.text = "Use cases: Market analysis, portfolio optimization, research reports"
    p.font.size = Pt(10)
    p.level = 1
    p.space_after = Pt(4)

    p = tf_left.add_paragraph()
    p.text = "Tools: get_fund_performance, get_market_data, run_portfolio_analysis"
    p.font.size = Pt(9)
    p.level = 1
    p.font.italic = True
    p.space_after = Pt(4)

    p = tf_left.add_paragraph()
    p.text = "Outcome: 40% faster research synthesis, compliance-aware"
    p.font.size = Pt(10)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Right column
    right_box = slide_c.shapes.add_textbox(
        Inches(5.5), Inches(1.3), Inches(right_width), Inches(column_height)
    )
    tf_right = right_box.text_frame
    tf_right.word_wrap = True

    # Compliance & Risk Agent
    p = tf_right.paragraphs[0]
    p.text = "Compliance & Risk Agent"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)

    p = tf_right.add_paragraph()
    p.text = "Use cases: Automated audits, regulatory impact analysis, risk scoring"
    p.font.size = Pt(10)
    p.level = 1
    p.space_after = Pt(4)

    p = tf_right.add_paragraph()
    p.text = "Tools: check_compliance, get_regulatory_requirements, trigger_workflow"
    p.font.size = Pt(9)
    p.level = 1
    p.font.italic = True
    p.space_after = Pt(4)

    p = tf_right.add_paragraph()
    p.text = "Outcome: 100% audit coverage, proactive risk identification"
    p.font.size = Pt(10)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_after = Pt(12)

    # HR Agent
    p = tf_right.add_paragraph()
    p.text = "HR Agent"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_before = Pt(8)

    p = tf_right.add_paragraph()
    p.text = "Use cases: Employee onboarding, policy Q&A, learning recommendations"
    p.font.size = Pt(10)
    p.level = 1
    p.space_after = Pt(4)

    p = tf_right.add_paragraph()
    p.text = "Tools: search_knowledge_base, get_employee_profile, trigger_workflow"
    p.font.size = Pt(9)
    p.level = 1
    p.font.italic = True
    p.space_after = Pt(4)

    p = tf_right.add_paragraph()
    p.text = "Outcome: 60% reduction in HR inquiries, faster onboarding"
    p.font.size = Pt(10)
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)

    # Add critical note at bottom
    note_box = slide_c.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
    )
    note_frame = note_box.text_frame
    note_para = note_frame.paragraphs[0]
    note_para.text = "Critical: All agents share the same core infrastructure. Each additional agent takes 3-4 weeks, not 6-12 months."
    note_para.font.size = Pt(11)
    note_para.font.bold = True
    note_para.font.color.rgb = RGBColor(192, 0, 0)
    note_para.alignment = PP_ALIGN.CENTER

    print("  ✓ Slide C created")

    # ========================================
    # SLIDE D: Platform Economics
    # ========================================
    print("Creating Slide D: The Economics of Platform Thinking...")
    slide_d = prs.slides.add_slide(blank_layout)

    # Add title
    create_text_box_with_title(slide_d, "The Economics of Platform Thinking", top=0.4)

    # Add content
    content_box = slide_d.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    # Traditional Approach
    p = tf.paragraphs[0]
    p.text = "Traditional Approach (Point Solutions):"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)
    p.space_after = Pt(8)

    bullets = [
        "Member Services AI: 12 months, $2M, separate infrastructure",
        "Investments AI: 12 months, $2M, separate infrastructure",
        "Compliance AI: 12 months, $2M, separate infrastructure",
        "HR AI: 12 months, $2M, separate infrastructure",
        "Total: 48 months (sequential), $8M, 4 separate systems to maintain"
    ]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(4)
        if "Total:" in bullet:
            p.font.bold = True

    # Platform Approach
    p = tf.add_paragraph()
    p.text = "Platform Approach:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_before = Pt(16)
    p.space_after = Pt(8)

    bullets = [
        "Foundation build: 3 months, $600K",
        "Member Services agent: 3 months, $500K (includes foundation)",
        "Investment agent: 4 weeks, $100K (reuses 70%)",
        "Compliance agent: 4 weeks, $100K (reuses 70%)",
        "HR agent: 4 weeks, $100K (reuses 70%)",
        "Total: 9 months (some parallel), $1.4M, one unified system"
    ]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(4)
        if "Total:" in bullet:
            p.font.bold = True

    # Savings
    p = tf.add_paragraph()
    p.text = "Savings: $6.6M (82%) and 75% faster time-to-value"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_before = Pt(16)
    p.alignment = PP_ALIGN.CENTER

    print("  ✓ Slide D created")

    # ========================================
    # SLIDE E: Ongoing Benefits
    # ========================================
    print("Creating Slide E: Platform Benefits...")
    slide_e = prs.slides.add_slide(blank_layout)

    # Add title
    create_text_box_with_title(slide_e, "Platform Benefits - Beyond Initial Savings", top=0.4)

    # Create two columns for benefits
    left_box = slide_e.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5)
    )
    tf_left = left_box.text_frame
    tf_left.word_wrap = True

    # Operational Benefits
    p = tf_left.paragraphs[0]
    p.text = "Operational Benefits:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_after = Pt(8)

    bullets = [
        "Single platform team vs 4 separate teams",
        "Shared monitoring and alerting",
        "One governance model to audit",
        "Centralized cost management"
    ]

    for bullet in bullets:
        p = tf_left.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(6)

    # Innovation Benefits
    p = tf_left.add_paragraph()
    p.text = "Innovation Benefits:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_before = Pt(16)
    p.space_after = Pt(8)

    bullets = [
        "New use case? Pick relevant tools from library",
        "Cross-agent workflows possible",
        "Learnings from one BU benefit all others",
        "Platform improves continuously for everyone"
    ]

    for bullet in bullets:
        p = tf_left.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(6)

    # Right column
    right_box = slide_e.shapes.add_textbox(
        Inches(5.5), Inches(1.3), Inches(4.5), Inches(5.5)
    )
    tf_right = right_box.text_frame
    tf_right.word_wrap = True

    # Risk Mitigation
    p = tf_right.paragraphs[0]
    p.text = "Risk Mitigation:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_after = Pt(8)

    bullets = [
        "Consistent security model",
        "No shadow AI deployments",
        "Complete audit trail across all BUs",
        "Regulatory confidence (single point of compliance)"
    ]

    for bullet in bullets:
        p = tf_right.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(6)

    # Strategic Advantage
    p = tf_right.add_paragraph()
    p.text = "Strategic Advantage:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    p.space_before = Pt(16)
    p.space_after = Pt(8)

    bullets = [
        "Most super funds build isolated solutions per BU",
        "ART will have unified intelligence across entire organization",
        "Faster response to market changes",
        "Better member experience through coordinated touchpoints"
    ]

    for bullet in bullets:
        p = tf_right.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(11)
        p.space_after = Pt(6)

    print("  ✓ Slide E created")

    # Save the presentation
    prs.save(output_file)
    print(f"\n{'='*70}")
    print(f"✅ SUCCESS! Presentation saved as: {output_file}")
    print(f"{'='*70}")
    print(f"\nAdded 4 new slides:")
    print("  1. Federated Multi-Agent Architecture")
    print("  2. Multi-BU Agent Deployment - Concrete Use Cases")
    print("  3. The Economics of Platform Thinking")
    print("  4. Platform Benefits - Beyond Initial Savings")
    print(f"\nTotal slides in presentation: {len(prs.slides)}")

if __name__ == "__main__":
    import sys

    # Check if python-pptx is installed
    try:
        from pptx import Presentation
    except ImportError:
        print("=" * 70)
        print("ERROR: python-pptx library is not installed.")
        print("=" * 70)
        print("Please install it using: pip install python-pptx")
        sys.exit(1)

    # Get file paths from command line if provided
    input_file = sys.argv[1] if len(sys.argv) > 1 else "ART_Enterprise_AI_Strategy.pptx"
    output_file = sys.argv[2] if len(sys.argv) > 2 else "ART_Enterprise_AI_Strategy_Enhanced.pptx"

    print("=" * 70)
    print("Enterprise Agentic Knowledge Platform - Slide Generator v2.0")
    print("=" * 70)
    print(f"Input file:  {input_file}")
    print(f"Output file: {output_file}")
    print("=" * 70)

    try:
        create_slides(input_file, output_file)
    except FileNotFoundError:
        print(f"\n❌ ERROR: Could not find input file: {input_file}")
        print("Please make sure the file exists and the path is correct.")
        sys.exit(1)
    except Exception as e:
        print(f"\n❌ ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
